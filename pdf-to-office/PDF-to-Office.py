#!/usr/bin/env python3
# -*- coding: utf-8 -*-
import os, threading, datetime, enum
from pathlib import Path
from tkinter import *
from tkinter import ttk, filedialog

from pdf2docx import Converter as Pdf2Doc
import pdfplumber
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt

# ---------------- PDF 转 PPT（可编辑文本） ----------------
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def pdf2pptx_editable_optimized(pdf_path, out_file):
    prs = Presentation()
    with pdfplumber.open(pdf_path) as pdf:
        for page_num, page in enumerate(pdf.pages, 1):
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
            text = page.extract_text()
            if not text:
                textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
                textbox.text_frame.add_paragraph().text = "(该页无文本)"
                continue

            lines = [line.strip() for line in text.split('\n') if line.strip()]
            y_offset = 0.5  # 初始纵向偏移
            for i, line in enumerate(lines):
                height = Inches(0.4)
                width = Inches(8)
                textbox = slide.shapes.add_textbox(Inches(0.5), Inches(y_offset), width, height)
                tf = textbox.text_frame
                tf.word_wrap = True
                p = tf.add_paragraph()
                p.text = line
                if i == 0:
                    # 第一行视作标题
                    p.font.size = Pt(24)
                    p.font.bold = True
                else:
                    p.font.size = Pt(18)
                p.alignment = PP_ALIGN.LEFT
                y_offset += 0.5  # 下一个段落下移
                if y_offset > 6:  # 超过幻灯片高度，停止添加
                    break
    prs.save(out_file)
# ---------------- PDF 转 Excel（纯 Python） ----------------
def pdf2excel(pdf_path, out_file):
    all_tables = []
    with pdfplumber.open(pdf_path) as pdf:
        for page_num, page in enumerate(pdf.pages, 1):
            tables = page.extract_tables()
            for table_index, table in enumerate(tables):
                if not table or len(table) < 2:
                    continue
                df = pd.DataFrame(table[1:], columns=table[0])
                all_tables.append((df, f"Page{page_num}_Table{table_index+1}"))

    if not all_tables:
        raise ValueError("PDF 中未找到可解析的表格")

    with pd.ExcelWriter(out_file, engine="openpyxl") as writer:
        for df, sheet_name in all_tables:
            df.to_excel(writer, sheet_name=sheet_name[:31], index=False)

# ---------------- 转换线程 ----------------
class ConvertThread(threading.Thread):
    def __init__(self, pdf_path, out_dir, fmt, progress_cb, msg_cb, done_cb):
        super().__init__(daemon=True)
        self.pdf_path = pdf_path
        self.out_dir = out_dir
        self.fmt = fmt
        self.progress_cb = progress_cb
        self.msg_cb = msg_cb
        self.done_cb = done_cb
        self._real_progress = 0
        self._fake_active = True

    def run(self):
        threading.Thread(target=self._fake_progress, daemon=True).start()
        try:
            out_file = pdf_convert(self.pdf_path, self.out_dir, self.fmt,
                                   self._update_real, self.msg_cb)
            self._fake_active = False
            self.progress_cb(100)
            self.done_cb(True, out_file)
        except Exception as e:
            self._fake_active = False
            self.done_cb(False, str(e))

    def _update_real(self, value):
        self._real_progress = value
        self.progress_cb(value)

    def _fake_progress(self):
        fake = 0
        while self._fake_active and fake < 99:
            if fake < self._real_progress:
                fake = self._real_progress
            else:
                fake += 1
            self.progress_cb(fake)
            threading.Event().wait(0.1)

# ---------------- 格式枚举 ----------------
class Format(enum.Enum):
    DOCX = ('Word 文档', '.docx')
    PPTX = ('PowerPoint', '.pptx')
    XLSX = ('Excel 表格', '.xlsx')

# ---------------- PDF 转 Office ----------------
def pdf_convert(pdf_path, out_dir, fmt: Format, progress_cb, msg_cb):
    out_dir = Path(out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    stamp = datetime.datetime.now().strftime("%Y%m%d%H%M%S")
    name = Path(pdf_path).stem
    out_file = str(out_dir / f"{name}_{stamp}{fmt.value[1]}")

    if fmt == Format.DOCX:
        cv = Pdf2Doc(pdf_path)
        total = len(cv.pages)
        cur = 0
        def cb(prog, desc):
            nonlocal cur
            if desc.get("event") == "page_parsed":
                cur += 1
                progress_cb(int(cur / total * 100))
                msg_cb(f"正在解析第 {cur}/{total} 页...")
        cv.convert(out_file, progress_callback=cb)
        cv.close()
    elif fmt == Format.PPTX:
        pdf2pptx_editable_optimized(pdf_path, out_file)
        progress_cb(100)
    elif fmt == Format.XLSX:
        pdf2excel(pdf_path, out_file)
        progress_cb(100)
    else:
        raise ValueError("不支持的目标格式")
    return out_file

# ---------------- GUI ----------------
class App(Tk):
    def __init__(self):
        super().__init__()
        self.title("PDF 转 Office 小工具")
        self.resizable(False, False)
        self.geometry("580x240")
        self.pdf_path = ""
        self.out_file = ""
        self.pdf_var = StringVar()
        self.out_var = StringVar(value=str(Path.home() / "Desktop/PDF转换结果"))
        self.fmt_var = StringVar(value='DOCX')

        pad = 12
        frm = Frame(self)
        frm.pack(fill=X, padx=pad, pady=pad)
        Label(frm, text="PDF 文件：", width=10, anchor=W).grid(row=0, column=0, sticky=W)
        Entry(frm, textvariable=self.pdf_var, width=45).grid(row=0, column=1, padx=(0,5))
        Button(frm, text="浏览...", width=8, command=self.browse_pdf).grid(row=0, column=2)

        Label(frm, text="输出目录：", width=10, anchor=W).grid(row=1, column=0, sticky=W, pady=(6,0))
        Entry(frm, textvariable=self.out_var, width=45).grid(row=1, column=1, padx=(0,5), pady=(6,0))
        Button(frm, text="浏览...", width=8, command=self.browse_out).grid(row=1, column=2, pady=(6,0))

        fmt_frm = Frame(self)
        fmt_frm.pack(fill=X, padx=pad, pady=(6,0))
        Label(fmt_frm, text="目标格式：", width=10, anchor=W).pack(side=LEFT)
        for fmt in Format:
            Radiobutton(fmt_frm, text=fmt.value[0], value=fmt.name,
                        variable=self.fmt_var).pack(side=LEFT, padx=8)

        self.progress = ttk.Progressbar(self, length=540, mode='determinate')
        self.progress.pack(pady=6, padx=pad)
        self.status = Label(self, text="准备就绪", anchor=W)
        self.status.pack(fill=X, padx=pad)

        btn_frm = Frame(self)
        btn_frm.pack(fill=X, padx=pad, pady=6)
        self.btn_start = Button(btn_frm, text="开始转换", width=12, command=self.start_convert)
        self.btn_start.pack(side=LEFT)
        self.btn_open = Button(btn_frm, text="打开文件", width=12, command=self.open_file)
        self.btn_open.pack_forget()

    def browse_pdf(self):
        path = filedialog.askopenfilename(title="选择 PDF 文件", filetypes=[("PDF 文件","*.pdf")])
        if path: self.pdf_var.set(path)

    def browse_out(self):
        path = filedialog.askdirectory(title="选择输出文件夹")
        if path: self.out_var.set(path)

    def start_convert(self):
        self.pdf_path = self.pdf_var.get()
        if not self.pdf_path:
            self.status.config(text="请先选择 PDF 文件！")
            return
        fmt = Format[self.fmt_var.get()]
        out_dir = self.out_var.get()
        if not out_dir:
            self.status.config(text="请选择输出目录！")
            return
        self.btn_open.pack_forget()
        self.status.config(text="正在转换，请稍候...")
        self.progress["value"] = 0
        self.btn_start.config(state=DISABLED)
        ConvertThread(
            self.pdf_path, out_dir, fmt,
            lambda v: self.progress.config(value=v),
            lambda s: self.status.config(text=s),
            self.on_finished
        ).start()

    def on_finished(self, ok: bool, out_file: str):
        self.btn_start.config(state=NORMAL)
        if ok:
            self.out_file = out_file
            self.status.config(text="转换成功！")
            self.btn_open.pack(side=RIGHT)
        else:
            self.status.config(text=f"转换失败：{out_file}")

    def open_file(self):
        if self.out_file and Path(self.out_file).is_file():
            os.startfile(self.out_file)

if __name__ == "__main__":
    App().mainloop()
